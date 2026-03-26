"""
Presentation Builder for EPL TOTW.

Creates a PL-styled presentation and exports to:
  - output/matchweek-{N}/presentation.pptx
  - output/matchweek-{N}/presentation.pdf

Reads:
  - output/matchweek-{N}/analysis/players.json
  - output/matchweek-{N}/analysis/formation.json
  - data/2025-26/matchweek-{N}/fixtures.json
  - data/2025-26/matchweek-{N+1}/fixtures.json  (next matchweek, optional)
  - output/matchweek-{N}/totw_diagram.png

CLI:
    python scripts/presentation_builder.py 21
"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import io
import json
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

sys.path.insert(0, str(Path(__file__).parent.parent))

import requests
from jinja2 import Environment, FileSystemLoader
from PIL import Image, ImageDraw

from scripts.utils import (
    load_json_cache,
    matchweek_analysis_dir,
    matchweek_data_dir,
    matchweek_output_dir,
    parse_fixture_date,
)

TEMPLATES_DIR = Path(__file__).parent.parent / "templates"

# ---------------------------------------------------------------------------
# Colors (for python-pptx PPTX building)
# ---------------------------------------------------------------------------
PURPLE     = (0x37, 0x00, 0x3C)
DARK_PURPLE= (0x2D, 0x00, 0x32)
GREEN      = (0x00, 0xFF, 0x87)
MAGENTA    = (0xE9, 0x00, 0x52)
WHITE      = (0xFF, 0xFF, 0xFF)
GREY       = (0xB3, 0xB3, 0xC6)
DARK_GREY  = (0x1A, 0x1A, 0x2E)

# ---------------------------------------------------------------------------
# Position full names
# ---------------------------------------------------------------------------
POSITION_FULL: dict[str, str] = {
    "GK":  "Goalkeeper",
    "RB":  "Right Back",
    "CB":  "Centre Back",
    "LB":  "Left Back",
    "RWB": "Right Wing Back",
    "LWB": "Left Wing Back",
    "CDM": "Def. Midfielder",
    "CM":  "Central Midfielder",
    "CAM": "Att. Midfielder",
    "RM":  "Right Midfielder",
    "LM":  "Left Midfielder",
    "RW":  "Right Winger",
    "LW":  "Left Winger",
    "CF":  "Centre Forward",
    "ST":  "Striker",
}


# ---------------------------------------------------------------------------
# Image download / cache
# ---------------------------------------------------------------------------

def _image_cache_dir(matchweek: int) -> Path:
    p = matchweek_output_dir(matchweek) / ".img_cache"
    p.mkdir(parents=True, exist_ok=True)
    return p


def _download_png(url: str, cache_dir: Path, timeout: int = 10) -> Optional[bytes]:
    """Download a PNG image and return bytes; None on failure."""
    if not url:
        return None
    key = hashlib.md5(url.encode()).hexdigest()
    cached = cache_dir / f"{key}.png"
    if cached.exists() and cached.stat().st_size > 100:
        return cached.read_bytes()
    try:
        resp = requests.get(url, timeout=timeout, headers={
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"
        })
        if resp.ok:
            data = resp.content
            cached.write_bytes(data)
            return data
    except Exception:
        pass
    return None


async def _svg_urls_to_png(urls: list[str], cache_dir: Path) -> dict[str, bytes]:
    """Batch-convert SVG URLs to PNG bytes using Playwright."""
    from playwright.async_api import async_playwright

    results: dict[str, bytes] = {}
    to_render = []
    for url in urls:
        if not url:
            continue
        key = hashlib.md5(url.encode()).hexdigest()
        cached = cache_dir / f"{key}.png"
        if cached.exists() and cached.stat().st_size > 100:
            results[url] = cached.read_bytes()
        else:
            to_render.append(url)

    if not to_render:
        return results

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        for url in to_render:
            try:
                page = await browser.new_page(viewport={"width": 100, "height": 100})
                await page.goto(url, wait_until="domcontentloaded", timeout=8000)
                await page.wait_for_timeout(500)
                png = await page.screenshot(full_page=True)
                await page.close()
                key = hashlib.md5(url.encode()).hexdigest()
                cached = cache_dir / f"{key}.png"
                cached.write_bytes(png)
                results[url] = png
            except Exception:
                pass
        await browser.close()

    return results


def _make_circular_png(img_bytes: bytes, size: int = 200) -> bytes:
    """Crop image to circle and return PNG bytes."""
    try:
        img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
        # Center-crop to square
        w, h = img.size
        s = min(w, h)
        img = img.crop(((w - s) // 2, (h - s) // 2, (w + s) // 2, (h + s) // 2))
        img = img.resize((size, size), Image.LANCZOS)
        # Apply circular mask
        mask = Image.new("L", (size, size), 0)
        ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
        result = Image.new("RGBA", (size, size), (0, 0, 0, 0))
        result.paste(img, mask=mask)
        out = io.BytesIO()
        result.save(out, format="PNG")
        return out.getvalue()
    except Exception:
        return img_bytes


# ---------------------------------------------------------------------------
# Stats extraction per position
# ---------------------------------------------------------------------------

def _v(d: dict, *keys, default=0):
    """Safely get a nested value from a dict."""
    val = d
    for k in keys:
        if not isinstance(val, dict):
            return default
        val = val.get(k, default)
    return val if val is not None else default


def _stat(label: str, value, highlight: bool = False) -> dict:
    return {"label": label, "value": str(value), "highlight": highlight}


def get_display_stats(position_slot: str, player: dict) -> list[dict]:
    """Return a list of stat dicts for display on the player slide."""
    s = player.get("stats", {})
    games    = s.get("games", {})
    goals    = s.get("goals", {})
    shots    = s.get("shots", {})
    passes   = s.get("passes", {})
    tackles  = s.get("tackles", {})
    duels    = s.get("duels", {})
    dribbles = s.get("dribbles", {})

    rating = games.get("rating") or "N/A"
    minutes = _v(games, "minutes", default=0)
    pos = position_slot.upper()

    if pos == "GK":
        saves = _v(goals, "saves", default=0)
        conceded = _v(goals, "conceded", default=0)
        clean = "Yes" if conceded == 0 else "No"
        return [
            _stat("SAVES", saves, highlight=True),
            _stat("GOALS CONCEDED", conceded),
            _stat("CLEAN SHEET", clean),
            _stat("MINUTES", minutes),
            _stat("MATCH RATING", rating, highlight=True),
        ]

    elif pos == "CB":
        return [
            _stat("TACKLES WON",     _v(tackles, "total",         default=0), highlight=True),
            _stat("INTERCEPTIONS",   _v(tackles, "interceptions",  default=0), highlight=True),
            _stat("CLEARANCES",      _v(tackles, "clearances",     default=0)),
            _stat("AERIAL DUELS WON",_v(duels,   "aerial_won",     default=0)),
            _stat("BLOCKS",          _v(tackles, "blocks",         default=0)),
            _stat("MATCH RATING",    rating, highlight=True),
        ]

    elif pos in ("RB", "LB", "RWB", "LWB"):
        def_actions = _v(tackles, "total", default=0) + _v(tackles, "interceptions", default=0)
        return [
            _stat("DEF. ACTIONS",   def_actions, highlight=True),
            _stat("KEY PASSES",      _v(passes,   "key",           default=0)),
            _stat("ASSISTS",         _v(goals,    "assists",        default=0)),
            _stat("TACKLES WON",     _v(tackles,  "total",         default=0)),
            _stat("INTERCEPTIONS",   _v(tackles,  "interceptions",  default=0)),
            _stat("MATCH RATING",    rating, highlight=True),
        ]

    elif pos == "CDM":
        return [
            _stat("TACKLES WON",    _v(tackles, "total",         default=0), highlight=True),
            _stat("INTERCEPTIONS",  _v(tackles, "interceptions",  default=0), highlight=True),
            _stat("CLEARANCES",     _v(tackles, "clearances",     default=0)),
            _stat("DUELS WON",      _v(duels,   "won",            default=0)),
            _stat("PASS ACCURACY",  f"{float(_v(passes, 'accuracy', default=0) or 0):.0f}%"),
            _stat("MATCH RATING",   rating, highlight=True),
        ]

    elif pos in ("CM", "CAM"):
        return [
            _stat("KEY PASSES",     _v(passes,   "key",      default=0), highlight=True),
            _stat("GOALS",          _v(goals,    "total",    default=0), highlight=True),
            _stat("ASSISTS",        _v(goals,    "assists",  default=0)),
            _stat("PASS ACCURACY",  f"{float(_v(passes, 'accuracy', default=0) or 0):.0f}%"),
            _stat("TACKLES WON",    _v(tackles,  "total",    default=0)),
            _stat("MATCH RATING",   rating, highlight=True),
        ]

    elif pos in ("RM", "LM", "RW", "LW"):
        return [
            _stat("GOALS",          _v(goals,    "total",    default=0), highlight=True),
            _stat("ASSISTS",        _v(goals,    "assists",  default=0), highlight=True),
            _stat("KEY PASSES",     _v(passes,   "key",      default=0)),
            _stat("DRIBBLES",       _v(dribbles, "success",  default=0)),
            _stat("SHOTS ON TARGET",_v(shots,    "on",       default=0)),
            _stat("MATCH RATING",   rating, highlight=True),
        ]

    elif pos in ("ST", "CF"):
        total_shots = _v(shots, "total", default=0)
        goals_val   = _v(goals, "total", default=0)
        conversion  = f"{round(goals_val / total_shots * 100)}%" if total_shots > 0 else "0%"
        return [
            _stat("GOALS",          goals_val, highlight=True),
            _stat("SHOTS ON TARGET",_v(shots, "on",       default=0), highlight=True),
            _stat("ASSISTS",        _v(goals,  "assists",  default=0)),
            _stat("SHOT CONVERSION",conversion),
            _stat("TOTAL SHOTS",    total_shots),
            _stat("MATCH RATING",   rating, highlight=True),
        ]

    # Fallback
    return [
        _stat("GOALS",   _v(goals, "total",   default=0), highlight=True),
        _stat("ASSISTS",  _v(goals, "assists", default=0)),
        _stat("MATCH RATING", rating, highlight=True),
    ]


# ---------------------------------------------------------------------------
# Data loading
# ---------------------------------------------------------------------------

def _fmt_date(date_str: str) -> str:
    """Format ISO date to short display string: 'Sat 3 Jan'."""
    try:
        dt = parse_fixture_date(date_str)
        return dt.strftime("%a %-d %b")
    except Exception:
        return date_str[:10]


def load_presentation_data(matchweek: int) -> dict:
    """Load all data needed for the presentation."""
    analysis_dir = matchweek_analysis_dir(matchweek)
    data_dir = matchweek_data_dir(matchweek)

    # Players
    players_data = load_json_cache(analysis_dir / "players.json")
    if not players_data:
        raise FileNotFoundError(f"No players.json for matchweek {matchweek}. Run player_evaluator.py first.")
    totw = players_data  # has 'players', 'matchweek', 'season', 'formation'

    # Formation
    formation_data = load_json_cache(analysis_dir / "formation.json")
    if not formation_data:
        raise FileNotFoundError(f"No formation.json for matchweek {matchweek}. Run formation_analyzer.py first.")

    # Fixtures (results)
    fixtures_raw = load_json_cache(data_dir / "fixtures.json") or []
    fixtures = []
    for f in fixtures_raw:
        fixtures.append({
            "home_name":  f["teams"]["home"]["name"],
            "home_logo":  f["teams"]["home"]["logo"],
            "home_score": f["goals"]["home"] if f["goals"]["home"] is not None else "–",
            "away_name":  f["teams"]["away"]["name"],
            "away_logo":  f["teams"]["away"]["logo"],
            "away_score": f["goals"]["away"] if f["goals"]["away"] is not None else "–",
            "date":       f["fixture"]["date"],
        })
    fixtures.sort(key=lambda x: x["date"])

    # Next matchweek fixtures
    next_mw = matchweek + 1
    next_data_dir = matchweek_data_dir(next_mw)
    next_fixtures_raw = load_json_cache(next_data_dir / "fixtures.json") or []
    next_fixtures = []
    for f in next_fixtures_raw:
        next_fixtures.append({
            "home_name": f["teams"]["home"]["name"],
            "home_logo": f["teams"]["home"]["logo"],
            "away_name": f["teams"]["away"]["name"],
            "away_logo": f["teams"]["away"]["logo"],
            "date_str":  _fmt_date(f["fixture"]["date"]),
        })
    next_fixtures.sort(key=lambda x: x["date_str"])

    # Build player list
    raw_players = totw.get("players", []) if isinstance(totw, dict) else totw
    formation_str = totw.get("formation", {}).get("formation", "4-3-3") if isinstance(totw, dict) else "4-3-3"

    players = []
    for tp in raw_players:
        p = tp["player"]
        slot = tp["position_slot"]
        stats = p.get("stats", {})
        initials = "".join(w[0].upper() for w in p["name"].split()[:2])

        players.append({
            "position_slot": slot,
            "position_full":  POSITION_FULL.get(slot, slot),
            "name":           p["name"],
            "initials":       initials,
            "team_name":      p.get("team_name", ""),
            "team_logo":      p.get("team_logo", ""),
            "photo_url":      p.get("photo", ""),
            "flag_url":       (f"https://media.api-sports.io/flags/{p['country_code']}.svg"
                               if p.get("country_code") and p["country_code"] not in ("xx", "")
                               else ""),
            "country_code":   p.get("country_code", ""),
            "key_stat":       tp.get("key_stat", ""),
            "selection_reason": tp.get("selection_reason", ""),
            "display_stats":  get_display_stats(slot, p),
            "stats":          stats,
        })

    # Formation details for slide 4
    formation_usages = formation_data.get("usages", [])

    # Winning teams (teams in the selected formation that won)
    selected_formation_usages = [u for u in formation_usages if u["formation"] == formation_str]
    winning_teams = []
    if selected_formation_usages:
        teams_in_formation = selected_formation_usages[0].get("teams", [])
        # Get logos from fixtures
        team_logo_map = {}
        for f in fixtures_raw:
            team_logo_map[f["teams"]["home"]["name"]] = f["teams"]["home"]["logo"]
            team_logo_map[f["teams"]["away"]["name"]] = f["teams"]["away"]["logo"]
        for tn in teams_in_formation:
            winning_teams.append({"name": tn, "logo": team_logo_map.get(tn, "")})

    return {
        "matchweek":          matchweek,
        "next_matchweek":     next_mw,
        "season":             totw.get("season", "2025/26") if isinstance(totw, dict) else "2025/26",
        "formation":          formation_str,
        "formation_rationale": formation_data.get("rationale", ""),
        "formation_usages":   formation_usages,
        "winning_teams":      winning_teams,
        "fixtures":           fixtures,
        "next_fixtures":      next_fixtures,
        "players":            players,
    }


# ---------------------------------------------------------------------------
# HTML → PDF via Playwright
# ---------------------------------------------------------------------------

def _diagram_data_uri(matchweek: int) -> str:
    """Load TOTW diagram PNG and return as base64 data URI."""
    png_path = matchweek_output_dir(matchweek) / "totw_diagram.png"
    if not png_path.exists():
        return ""
    b64 = base64.b64encode(png_path.read_bytes()).decode()
    return f"data:image/png;base64,{b64}"


def render_html(data: dict) -> str:
    """Render the presentation.html Jinja2 template."""
    env = Environment(loader=FileSystemLoader(str(TEMPLATES_DIR)))
    template = env.get_template("presentation.html")
    return template.render(**data)


async def _export_pdf(html_path: Path, pdf_path: Path) -> None:
    """Use Playwright to print the HTML presentation to a PDF."""
    from playwright.async_api import async_playwright

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        page = await browser.new_page(viewport={"width": 1280, "height": 720})
        await page.goto(f"file://{html_path.resolve()}", wait_until="networkidle")
        # Wait for all images to load
        await page.wait_for_timeout(4000)
        await page.pdf(
            path=str(pdf_path),
            width="1280px",
            height="720px",
            print_background=True,
            margin={"top": "0", "bottom": "0", "left": "0", "right": "0"},
        )
        await browser.close()


def build_pdf(matchweek: int, data: dict, output_dir: Path) -> Path:
    """Render HTML presentation and export to PDF."""
    data_with_diagram = {**data, "diagram_data_uri": _diagram_data_uri(matchweek)}
    html = render_html(data_with_diagram)

    html_path = output_dir / "presentation.html"
    html_path.write_text(html, encoding="utf-8")
    print(f"  HTML saved: {html_path}")

    pdf_path = output_dir / "presentation.pdf"
    asyncio.run(_export_pdf(html_path, pdf_path))
    size_kb = pdf_path.stat().st_size // 1024
    print(f"  PDF saved:  {pdf_path} ({size_kb} KB)")
    return pdf_path


# ---------------------------------------------------------------------------
# PPTX builder (python-pptx)
# ---------------------------------------------------------------------------

def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _fill(shape, r, g, b):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(r, g, b)


def _add_text(slide, text: str, left, top, width, height,
              font_size: int, font_color, bold: bool = False, align="left",
              word_wrap: bool = True) -> None:
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(*font_color)
    run.font.bold = bold
    run.font.name = "Arial"


def _set_bg(slide, r, g, b):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(r, g, b)


def _add_img_bytes(slide, img_bytes: bytes, left, top, width, height) -> None:
    from pptx.util import Inches
    try:
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)
    except Exception:
        pass


def build_pptx(matchweek: int, data: dict, output_dir: Path) -> Path:
    """Build the PPTX presentation using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height
    blank_layout = prs.slide_layouts[6]  # blank

    img_cache = _image_cache_dir(matchweek)

    def px(val_inches: float) -> Emu:
        return Inches(val_inches)

    # ---- Slide 1: Title ----
    s = prs.slides.add_slide(blank_layout)
    _set_bg(s, *PURPLE)
    _add_text(s, "PREMIER LEAGUE · 2025/26", px(0), px(2.8), W, px(0.4),
              14, GREEN, bold=True, align="center")
    _add_text(s, "TEAM OF THE WEEK", px(0), px(3.3), W, px(1.4),
              54, WHITE, bold=True, align="center")
    _add_text(s, f"MATCHWEEK {matchweek}", px(0), px(5.0), W, px(0.7),
              28, GREEN, bold=True, align="center")
    _add_text(s, f"{data['season']} SEASON", px(0), px(5.8), W, px(0.4),
              14, GREY, align="center")

    # ---- Slide 2: Results ----
    s = prs.slides.add_slide(blank_layout)
    _set_bg(s, *PURPLE)
    _add_text(s, f"MATCHWEEK {matchweek} · RESULTS", px(0.5), px(0.3), W, px(0.6),
              24, WHITE, bold=True)
    col_w = px(6.0)
    col_h = px(0.68)
    for i, f in enumerate(data["fixtures"]):
        col = i // 5
        row = i % 5
        lft = px(0.4 + col * 6.5)
        top = px(1.1 + row * 0.76)
        score = f"{f['home_score']}–{f['away_score']}"
        line = f"{f['home_name'][:15]:<15}  {score}  {f['away_name'][:15]}"
        _add_text(s, line, lft, top, col_w, col_h, 13, WHITE)

    # ---- Slide 3: Section Formation ----
    s = prs.slides.add_slide(blank_layout)
    _set_bg(s, *DARK_PURPLE)
    _add_text(s, "TOTW  FORMATION", px(0), px(2.5), W, px(0.6), 18, GREY, align="center")
    _add_text(s, data["formation"], px(0), px(3.2), W, px(1.8), 96, WHITE, bold=True, align="center")

    # ---- Slide 4: Formation Explanation ----
    s = prs.slides.add_slide(blank_layout)
    _set_bg(s, *PURPLE)
    _add_text(s, data["formation"], px(0.5), px(0.3), px(3), px(1.6), 72, GREEN, bold=True)
    _add_text(s, "SELECTED FORMATION", px(0.5), px(0.2), px(4), px(0.4), 12, GREY, bold=True)
    _add_text(s, data["formation_rationale"], px(0.5), px(2.0), px(5.5), px(3.0), 14, GREY)
    _add_text(s, "FORMATION USAGE THIS WEEK", px(6.8), px(0.3), px(6), px(0.4), 12, GREY, bold=True)
    for i, u in enumerate(data["formation_usages"][:6]):
        top = px(0.8 + i * 0.88)
        line = f"{u['formation']}  |  {u['win_count']}W / {u['usage_count']}G  |  {', '.join(u['teams'][:2]) or '–'}"
        _add_text(s, line, px(6.8), top, px(6.0), px(0.5), 14, WHITE)

    # ---- Slide 5: Section Players ----
    s = prs.slides.add_slide(blank_layout)
    _set_bg(s, *DARK_PURPLE)
    _add_text(s, "TOTW  PLAYERS", px(0), px(2.5), W, px(0.6), 18, GREY, align="center")
    _add_text(s, "XI SELECTED", px(0), px(3.3), W, px(1.4), 60, WHITE, bold=True, align="center")

    # ---- Slides 6–16: Player slides ----
    for player in data["players"]:
        s = prs.slides.add_slide(blank_layout)
        _set_bg(s, *DARK_PURPLE)

        # Left panel background
        rect = s.shapes.add_shape(1, px(0), px(0), px(4.2), H)
        _fill(rect, *PURPLE)

        # Player photo (circular, if available)
        photo_bytes = _download_png(player["photo_url"], img_cache)
        if photo_bytes:
            try:
                circ = _make_circular_png(photo_bytes, size=240)
                _add_img_bytes(s, circ, px(0.7), px(0.9), px(2.8), px(2.8))
            except Exception:
                pass

        # Team badge
        badge_url = player["team_logo"]
        if badge_url:
            badge_bytes = _download_png(badge_url, img_cache)
            if not badge_bytes and badge_url.endswith(".svg"):
                # Try PNG variant
                png_url = badge_url.replace("/rb/", "/50/").replace(".svg", "@x2.png")
                badge_bytes = _download_png(png_url, img_cache)
            if badge_bytes:
                try:
                    _add_img_bytes(s, badge_bytes, px(2.8), px(3.0), px(0.9), px(0.9))
                except Exception:
                    pass

        # Player name + details on left
        _add_text(s, player["name"], px(0.2), px(4.0), px(3.8), px(0.8), 18, WHITE, bold=True, align="center")
        _add_text(s, player["team_name"], px(0.2), px(4.8), px(3.8), px(0.4), 12, GREY, align="center")
        pos_lbl_top = px(5.2)
        _add_text(s, player["position_slot"], px(0.2), pos_lbl_top, px(3.8), px(0.5), 14, GREEN, bold=True, align="center")

        # Right panel — stats
        x_right = px(4.6)
        _add_text(s, player["position_full"].upper(), x_right, px(0.3), px(8.4), px(0.8), 36, MAGENTA, bold=True)
        _add_text(s, player["name"], x_right, px(1.1), px(8.4), px(0.5), 18, WHITE, bold=True)
        _add_text(s, player["key_stat"], x_right, px(1.6), px(8.4), px(0.4), 12, GREEN)

        # Stats grid (2 cols × 3 rows)
        for idx, stat in enumerate(player["display_stats"][:6]):
            col = idx % 2
            row = idx // 2
            sx = x_right + px(col * 4.1)
            sy = px(2.3 + row * 1.6)
            _add_text(s, stat["label"], sx, sy, px(3.8), px(0.3), 10, GREY, bold=True)
            color = GREEN if stat["highlight"] else WHITE
            _add_text(s, stat["value"], sx, sy + px(0.3), px(3.8), px(0.9), 32, color, bold=True)

    # ---- Slide 17: Diagram + Next Fixtures ----
    s = prs.slides.add_slide(blank_layout)
    _set_bg(s, *PURPLE)

    # Diagram image (left half)
    diagram_path = matchweek_output_dir(matchweek) / "totw_diagram.png"
    if diagram_path.exists():
        try:
            s.shapes.add_picture(str(diagram_path), px(0), px(0), px(6.2), H)
        except Exception:
            pass

    # Right panel
    right_rect = s.shapes.add_shape(1, px(6.2), px(0), px(7.1), H)
    _fill(right_rect, *DARK_PURPLE)
    _add_text(s, "COMING UP", px(6.6), px(0.5), px(6.5), px(0.4), 13, GREY, bold=True)
    _add_text(s, f"MATCHWEEK {data['next_matchweek']}", px(6.6), px(1.0), px(6.5), px(0.7), 28, GREEN, bold=True)

    if data["next_fixtures"]:
        for i, nf in enumerate(data["next_fixtures"][:10]):
            top = px(1.9 + i * 0.56)
            line = f"{nf['home_name'][:14]:<14}  vs  {nf['away_name'][:14]:<14}  {nf['date_str']}"
            _add_text(s, line, px(6.6), top, px(6.5), px(0.5), 12, WHITE)
    else:
        _add_text(s, "Next matchweek fixtures not yet available.", px(6.6), px(2.2), px(6.5), px(0.5), 13, GREY)

    pptx_path = output_dir / "presentation.pptx"
    prs.save(str(pptx_path))
    size_kb = pptx_path.stat().st_size // 1024
    print(f"  PPTX saved: {pptx_path} ({size_kb} KB)")
    return pptx_path


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def build_presentation(matchweek: int) -> tuple[Path, Path]:
    """Full pipeline: load data → build HTML/PDF + PPTX."""
    print(f"\nBuilding presentation for matchweek {matchweek}...")

    data = load_presentation_data(matchweek)
    print(f"  Formation:  {data['formation']}")
    print(f"  Players:    {len(data['players'])}")
    print(f"  Fixtures:   {len(data['fixtures'])}")
    print(f"  Next MW:    {data['next_matchweek']} ({len(data['next_fixtures'])} fixtures)")

    output_dir = matchweek_output_dir(matchweek)

    print("\n→ Building PDF (HTML + Playwright)...")
    pdf_path = build_pdf(matchweek, data, output_dir)

    print("\n→ Building PPTX (python-pptx)...")
    pptx_path = build_pptx(matchweek, data, output_dir)

    print(f"\nDone.")
    print(f"  PDF:  {pdf_path}")
    print(f"  PPTX: {pptx_path}")
    return pdf_path, pptx_path


def main():
    if len(sys.argv) < 2:
        print("Usage: python scripts/presentation_builder.py <matchweek>")
        sys.exit(1)
    matchweek = int(sys.argv[1])
    build_presentation(matchweek)


if __name__ == "__main__":
    main()
